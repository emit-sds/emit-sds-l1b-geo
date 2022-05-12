# Stuff we can generally find useful
import geocal
import emit
import glob
import os
import multiprocessing
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

def ppt_save(ppt_file):
    '''Write out the latest matplotlib data to the powerpoint file'''
    blank_slide_layout = ppt_file.slide_layouts[6]
    slide = ppt_file.slides.add_slide(blank_slide_layout)
    img = BytesIO()
    plt.tight_layout()
    plt.savefig(img, dpi=300)
    pic = slide.shapes.add_picture(img, Inches(0), Inches(0), height=Inches(7.5))

def plot_camera_data(gcam, title):
    '''Plot the camera data out.'''
    plt.clf()
    fig, ax = plt.subplots(2,1)
    fig.suptitle(title)
    ax[0].plot(list(range(gcam.field_alignment.shape[0])),
            gcam.field_alignment[:,0])
    ax[0].set(xlabel="Camera Sample", ylabel="Field Alignment X",
              title="Camera field aligment X")
    ax[1].plot(list(range(gcam.field_alignment.shape[0])),
          gcam.field_alignment[:,1])
    ax[1].set(xlabel="Camera Sample", ylabel="Field Alignment Y",
              title="Camera field aligment Y")
    plt.tight_layout()
    
def find_exactly_one_file(pattern):
    f = glob.glob(pattern)
    if(len(f) == 0):
        raise RuntimeError(f"Could find a file matching {pattern}")
    if(len(f) > 1):
        raise RuntimeError(f"Found more than one file matching {pattern}")
    return f[0]

def create_igc(dir, with_rad = True, pyortho_loc = False):
    '''Create a IGC that has the final results for a given run directory.
    We have an assumed set of files in the directory, just to make this
    simpler.'''
    # Strip off trailing "/" if it is supplied
    if(os.path.basename(dir) == ''):
        dir = os.path.dirname(dir)
    bname = os.path.basename(dir).split("_")[0]
    orb = emit.AvirisNgOrbit(find_exactly_one_file(dir + "/*_att.nc"))
    tt = emit.AvirisNgTimeTable(find_exactly_one_file(dir + "/*_line_time.nc"))
    rad_band = 54
    rad_scale = 100.0
    if(with_rad):
        rad = geocal.GdalRasterImage(find_exactly_one_file(dir + "/*_rdn_v2z4_clip"),
                                     rad_band)
        rad = geocal.ScaleImage(rad, rad_scale)
    else:
        rad = None
    if(os.path.exists(f"{dir}/{bname}_camera.xml")):
       cam = geocal.read_shelve(f"{dir}/{bname}_camera.xml")
    else:
       l1_osp_dir = emit.L1OspDir("l1_osp")
       cam = l1_osp_dir.camera() 
    ipi = geocal.Ipi(orb,cam,0,tt.min_time, tt.max_time,tt)
    igc = geocal.IpiImageGroundConnection(ipi, geocal.SrtmDem(), rad,
                                          bname, 1.0)
    # Tack on the loc file, it is useful to have access to this
    if(os.path.exists(f"{dir}/{bname}_loc")):
        igc.loc = emit.AvirisNgLoc(f"{dir}/{bname}_loc")
    else:
        igc.loc = None
    if(pyortho_loc):
        if(os.path.exists(f"/beegfs/store/ang/y22/raw/{bname}_rdn_loc")):
            igc.pyortho_loc = emit.AvirisNgLoc(f"/beegfs/store/ang/y22/raw/{bname}_rdn_loc")
        else:
            igc.pyortho_loc = None
    return igc

def dir_igc(igc):
    '''Return the directory that the igc is for.'''
    return os.path.dirname(igc.image.raw_data.file_names[0])
    
def rot_mi(igc):
    mi = geocal.cib01_mapinfo(igc.resolution_meter())
    loc = igc.loc
    return loc.map_info_rotated(mi)

        
