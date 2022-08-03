# Short script to put the pyortho camera model into a camera model
# we can use in geocal.

from emit import *
from geocal import *
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

def ppt_save(ppt_file):
    blank_slide_layout = ppt_file.slide_layouts[6]
    slide = ppt_file.slides.add_slide(blank_slide_layout)
    img = BytesIO()
    plt.tight_layout()
    plt.savefig(img, dpi=300)
    pic = slide.shapes.add_picture(img, Inches(0), Inches(0), height=Inches(7.5))
    
def check_cam(gcam, camortho):
    '''Do a simple comparison between the GLAS camera model and the original
    camortho data.'''
    ldiff = []
    sdiff = []
    for smp in range(gcam.number_sample(0)):
        slv = ScLookVector(camortho[smp,0], camortho[smp,1], camortho[smp,2])
        fc = gcam.frame_coordinate(slv,0)
        ldiff.append(fc.line-0)
        sdiff.append(fc.sample-smp)
    print("Line difference ", abs(np.array(ldiff)).max())
    print("Sample difference ", abs(np.array(sdiff)).max())

def plot_data(gcam, title):
    '''Plot the camera data out.'''
    plt.clf()
    fig, ax = plt.subplots(2,1)
    fig.suptitle(title)
    ax[0].plot(list(range(gcam.field_alignment.shape[0])),
            gcam.field_alignment[:,0])
    ax[0].set(xlabel="Camera Sample", ylabel="Field Alignment X",
              title="Camera field aligment X")
    ax[1].plot(list(range(gcam.field_alignment.shape[0])),
          gcam.field_alignment[:,1])
    ax[1].set(xlabel="Camera Sample", ylabel="Field Alignment Y",
              title="Camera field aligment Y")
    plt.tight_layout()

ppt = Presentation()        
for f in ("avng_2014_camera_cal3_2014_06_20_14_26_06",
          "avng_b200_2017_camera_cal3_2017_03_23_15_19_17",
          "avng_er2_2017_camera_cal3_2017_03_23_15_19_17"):
    # This is the ScLookVector for each sample.
    camortho = np.fromfile(f, dtype='<d').reshape([-1,3])
    sclv_list = [ScLookVector(camortho[smp,0], camortho[smp,1], camortho[smp,2])
                 for smp in range(camortho.shape[0])]
    gcam = GlasGfmCamera.create_glas_from_sc_look_vector(sclv_list,
                                                         focal_length=27.5)
    write_shelve(f"{f}_glas.xml", gcam)
    check_cam(gcam, camortho)
    plot_data(gcam, f"{f} Field Alignment")
    ppt_save(ppt)
ppt.save("camera_model.ppt")
