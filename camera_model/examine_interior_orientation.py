from geocal import *
from emit import *
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

dname = "/home/smyth/LocalFast/Emit/orbit_22204/output/"
igccol = GeoQa.igccol_corrected(dname + "emit20220810t051018_o22204_l1b_geoqa_b001_v01.nc")
# Selected two scenes that have a good number of tiepoints collected.
# Data looks pretty good visually also
igc1 = igccol.image_ground_connection(18)
# Note the index here is 1 based, to we want 19 instead of 18
img1 = GdalRasterImage(dname + "proj_erdas_019.img", 1)
rimg1 = GdalRasterImage(dname + "ref_erdas_019.img", 1)
rigc1 = MapInfoImageGroundConnection(rimg1, igc1.dem)
igc2 = igccol.image_ground_connection(19)
img2 = GdalRasterImage(dname + "proj_erdas_020.img", 1)
rimg2 = GdalRasterImage(dname + "ref_erdas_020.img", 1)
rigc2 = MapInfoImageGroundConnection(rimg2, igc2.dem)
m = CcorrLsmMatcher()
sm2 = SurfaceImageToImageMatch(igc2,img2,rigc2,rimg2, m)

def collinearity_residual(parm, igc, tpcol):
    '''Simple calculation of collinearity residual. This is a simpler
       interface than doing a full SBA, and should be good for doing
       an initial camera exterior orientation calculation.
    '''
    igc.ipi.camera.full_camera.parameter_subset = parm
    res = []
    for tp in tpcol:
        gp = tp.ground_location
        if(tp.image_coordinate(0)):
            res.extend(igc.collinearity_residual(gp, tp.image_coordinate(0)))
    return res

def fit_data(igc, tpcol):
    cfull = igc.ipi.camera.full_camera
    cfull.fit_epsilon = True
    cfull.fit_beta = True
    cfull.fit_delta = True
    cfull.fit_focal_length = False

    x0 = cfull.parameter_subset.copy()
    r = scipy.optimize.least_squares(collinearity_residual, x0,
                                     args=(igc, tpcol), loss = "huber")
    print(r)

igclist = [igc1, igc2]
rigclist = [rigc1, rigc2]
imglist = [img1, img2]
rimglist = [rimg1, rimg2]
dflist = []
for i in range(len(igclist)):
    igc = igclist[i]
    rigc = rigclist[i]
    sm = SurfaceImageToImageMatch(igc,imglist[i],rigc,rimglist[i], m)
    tpcol = []
    for ln in range(10,igc.number_line-10,10):
        print(f"Starting line {ln}. Have found {len(tpcol)} matches")
        for smp in range(i % 10,igc1.number_sample-10,10):
            ic1,ic2,_,_,success,_ = sm.match_surf(igc.ground_coordinate(ImageCoordinate(ln,smp)))
            if(success):
                tp = TiePoint(1)
                tp.is_gcp = True
                tp.ground_location = rigc.ground_coordinate(ic2)
                tp.image_coordinate(0, ic1)
                tpcol.append(tp)
    print(f"Found a total of  {len(tpcol)} tiepoints, image {igc.number_line} x {igc.number_sample}")
    tpcol = TiePointCollection(tpcol)
    # Looking just for pixel to pixel differences, so fit out any
    # average error by making a small tweak to the exterior orientation
    fit_data(igc, tpcol)
    res = []
    for tp in tpcol:
        ic1pred = igc.image_coordinate(tp.ground_location)
        ic1 = tp.image_coordinate(0)
        res.append([ic1.line, ic1.sample, ic1.line - ic1pred.line, ic1.sample-ic1pred.sample])
    df = pd.DataFrame(np.array(res), columns=["Line", "Sample", "Delta Line",
                                              "Delta Sample"])
    dflist.append(df)

def ppt_save(ppt_file):
    '''Write out the latest matplotlib data to the powerpoint file'''
    blank_slide_layout = ppt_file.slide_layouts[6]
    slide = ppt_file.slides.add_slide(blank_slide_layout)
    img = BytesIO()
    plt.tight_layout()
    plt.savefig(img, dpi=300)
    pic = slide.shapes.add_picture(img, Inches(0), Inches(0), height=Inches(7.5))

def plot_delta_data(df, title, ylim = None):
    '''Plot the camera data out.'''
    plt.clf()
    fig, ax = plt.subplots(2,1)
    fig.suptitle(title)
    df.plot("Sample","Delta Sample", ax=ax[0], ylim=ylim,
             kind="hexbin",
             title="Density of Difference Sample Predicted vs. Image Match")
    df.plot("Sample","Delta Line", ylim=ylim, ax=ax[1], kind="hexbin",
            title="Density of Difference Line Predicted vs. Image Match")
    plt.tight_layout()

# We should probably allow a tweaking of the fit here based on tiepoints.
#df.plot("Sample","Delta Line", kind="hexbin")
ppt = Presentation()
plot_delta_data(dflist[0], f"Interior Orientation {igccol.scene_time_list[18]}")
ppt_save(ppt)
plot_delta_data(dflist[0], f"Close up {igccol.scene_time_list[18]}", ylim=[-2,2])
ppt_save(ppt)
plot_delta_data(dflist[1], f"Interior Orientation {igccol.scene_time_list[19]}")
ppt_save(ppt)
plot_delta_data(dflist[1], f"Close up {igccol.scene_time_list[19]}", ylim=[-2,2])
ppt_save(ppt)
ppt.save("interior_orientation_plot.ppt")
